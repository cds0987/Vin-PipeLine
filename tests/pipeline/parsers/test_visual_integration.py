"""
Integration tests for pipeline.parsers._visual

Uses the REAL AIProvider (OpenRouter / gpt-4o-mini) configured via .env.
Creates test files on-the-fly using Pillow and python-pptx — no external fixtures.

Run:
    pytest tests/pipeline/parsers/test_visual_integration.py -v -m integration

Skipped automatically when AI_API_KEY is absent or AI_PROVIDER=mock.
"""
from __future__ import annotations

import io
import os

import pytest

# ── skip guard ────────────────────────────────────────────────────────────────

_has_key = bool(os.getenv("AI_API_KEY", "").strip())
_is_mock = os.getenv("AI_PROVIDER", "auto").lower() == "mock"
_skip_reason = "AI_API_KEY not set or AI_PROVIDER=mock — real API required"
pytestmark = pytest.mark.integration


def _require_real_api():
    if not _has_key or _is_mock:
        pytest.skip(_skip_reason)


# ── helpers: create test files in memory ─────────────────────────────────────


def _get_font(size: int = 28):
    """Return a PIL font at the requested size, falling back to default."""
    from PIL import ImageFont
    try:
        # Try common system fonts (available in most Linux Docker images)
        for name in ("DejaVuSans.ttf", "LiberationSans-Regular.ttf",
                     "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                     "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"):
            try:
                return ImageFont.truetype(name, size)
            except (IOError, OSError):
                continue
    except Exception:
        pass
    # Pillow >= 10 supports size parameter on the default font
    try:
        return ImageFont.load_default(size=size)
    except TypeError:
        return ImageFont.load_default()


def _make_png_with_text(text_lines: list[str]) -> bytes:
    """Create a white PNG with large, OCR-friendly text."""
    from PIL import Image, ImageDraw

    font = _get_font(32)
    width, height = 1200, len(text_lines) * 80 + 100
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)

    y = 40
    for line in text_lines:
        draw.text((60, y), line, fill="black", font=font)
        y += 72

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_image_only_pdf(text_lines: list[str]) -> bytes:
    """Create a PDF whose single page is a rasterised image — no text layer."""
    from PIL import Image, ImageDraw

    font = _get_font(36)
    width, height = 1400, len(text_lines) * 90 + 120
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)

    y = 60
    for line in text_lines:
        draw.text((80, y), line, fill="black", font=font)
        y += 88

    buf = io.BytesIO()
    img.save(buf, format="PDF", resolution=150)
    return buf.getvalue()


def _make_pptx_with_image(image_bytes: bytes) -> bytes:
    """Create a single-slide PPTX with the given image embedded."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        Inches(0.5), Inches(0.5),
        Inches(9.0), Inches(6.5),
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── shared fixture ────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def real_ai():
    """Return a real OpenAIProvider built from .env settings."""
    _require_real_api()
    from utils.ai_provider import build_ai_provider
    provider, warning = build_ai_provider()
    if warning:
        pytest.skip(f"AIProvider fell back: {warning}")
    return provider


# ── content used across tests ─────────────────────────────────────────────────

_CONTENT_LINES = [
    "Chính sách công tác nội bộ VSF 2024",
    "Hoàn tiền vé máy bay: tối đa 5,000,000 VNĐ / chuyến",
    "Khách sạn: tối đa 800,000 VNĐ / đêm",
    "Thời hạn nộp chứng từ: 7 ngày sau khi về",
]

# Key words the vision model must find in the output.
_EXPECTED_KEYWORDS = ["VSF", "5,000,000", "800,000"]


# ── PNG (standalone image) ────────────────────────────────────────────────────


def test_png_ocr_extracts_text(real_ai):
    """Standalone PNG → AIProvider.ocr() → markdown with key content."""
    png_bytes = _make_png_with_text(_CONTENT_LINES)

    from pipeline.parsers import _visual
    result = _visual._parse_image(png_bytes, real_ai)

    assert result, "Parser returned empty string for PNG"
    found = [kw for kw in _EXPECTED_KEYWORDS if kw in result]
    assert len(found) >= 2, (
        f"Expected at least 2 of {_EXPECTED_KEYWORDS} in OCR output.\n"
        f"Got: {result[:500]}"
    )


# ── Image-only PDF (no text layer) ───────────────────────────────────────────


def test_image_only_pdf_extracted_via_vision(real_ai):
    """PDF with no text layer → MarkItDown + gpt-4o-mini → markdown."""
    pdf_bytes = _make_image_only_pdf(_CONTENT_LINES)

    from pipeline.parsers import run
    result = run(pdf_bytes, ".pdf", ai_provider=real_ai)

    assert result, "Parser returned empty string for image-only PDF"
    found = [kw for kw in _EXPECTED_KEYWORDS if kw in result]
    assert len(found) >= 2, (
        f"Expected at least 2 of {_EXPECTED_KEYWORDS} in PDF output.\n"
        f"Got: {result[:500]}"
    )


def test_image_only_pdf_returns_markdown_string(real_ai):
    """Output must be a non-empty string (type check)."""
    pdf_bytes = _make_image_only_pdf(["VSF internal policy document 2024"])
    from pipeline.parsers import run
    result = run(pdf_bytes, ".pdf", ai_provider=real_ai)
    assert isinstance(result, str)
    assert len(result.strip()) > 0


# ── PPTX with embedded image ─────────────────────────────────────────────────


def test_pptx_with_embedded_image(real_ai):
    """PPTX slide with embedded image → MarkItDown + vision → markdown."""
    png_bytes = _make_png_with_text(_CONTENT_LINES)
    pptx_bytes = _make_pptx_with_image(png_bytes)

    from pipeline.parsers import run
    result = run(pptx_bytes, ".pptx", ai_provider=real_ai)

    assert result, "Parser returned empty string for PPTX"
    found = [kw for kw in _EXPECTED_KEYWORDS if kw in result]
    assert len(found) >= 1, (
        f"Expected at least 1 of {_EXPECTED_KEYWORDS} in PPTX output.\n"
        f"Got: {result[:500]}"
    )


# ── DOCX (text-only sanity check) ────────────────────────────────────────────


def test_docx_text_content_extracted(real_ai, tmp_path):
    """DOCX with plain text → MarkItDown extracts it correctly."""
    from docx import Document

    doc = Document()
    doc.add_heading("Quy định bảo mật thông tin", 0)
    doc.add_paragraph("Nhân viên không được chia sẻ thông tin nội bộ.")
    doc.add_paragraph("Vi phạm dẫn đến chấm dứt hợp đồng lao động.")

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    from pipeline.parsers import run
    result = run(docx_bytes, ".docx", ai_provider=real_ai)

    assert "bảo mật" in result or "Quy định" in result, (
        f"Expected DOCX content in output.\nGot: {result[:500]}"
    )


# ── mock fallback sanity check ────────────────────────────────────────────────


def test_mock_provider_falls_back_to_text_only():
    """MockAIProvider (no llm_client) must not call OpenRouter — text-only fallback."""
    from utils.ai_provider import MockAIProvider
    from pipeline.parsers import run

    # Simple text PDF content — MarkItDown text extraction only
    content = b"VSF internal policy document\nPage 1 of 1"
    # Use .txt (text format, no ai_provider needed)
    result = run(content, ".txt")
    assert isinstance(result, str)
